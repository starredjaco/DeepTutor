# -*- coding: utf-8 -*-
"""
Shared multimodal utilities for benchmark evaluation.

Provides a unified pipeline for processing file attachments (images,
documents, spreadsheets, text) and building OpenAI-compatible multimodal
messages.  Used by GAIA, HLE, and any future benchmark that may include
multimodal questions.

Design principles:
  - If the configured model supports multimodal input, images are encoded
    as base64 data URIs and sent via the ``image_url`` content type.
  - If the model does NOT support multimodal, images are described as
    text placeholders so the run still proceeds (text-only mode).
  - Non-image files (xlsx, csv, pdf, docx, txt, py, …) are always read
    as text and injected directly into the prompt.
  - Every attachment processing step is logged so the user can verify
    exactly what was sent to the LLM.
"""

from __future__ import annotations

import base64
import logging
import mimetypes
import os
from dataclasses import dataclass, field
from typing import Any

logger = logging.getLogger("benchmark.multimodal")

# ------------------------------------------------------------------
# Supported file categories
# ------------------------------------------------------------------

IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"})
TEXT_EXTENSIONS = frozenset({".txt", ".py", ".csv", ".json", ".jsonld", ".md", ".html", ".xml", ".yaml", ".yml"})
SPREADSHEET_EXTENSIONS = frozenset({".xlsx", ".xls"})
DOCUMENT_EXTENSIONS = frozenset({".pdf", ".docx", ".pptx"})

MAX_TEXT_CHARS = 50_000


# ------------------------------------------------------------------
# Data structures
# ------------------------------------------------------------------

@dataclass
class AttachmentResult:
    """Result of processing a single file attachment."""
    file_path: str
    file_name: str
    file_type: str          # e.g. "image", "spreadsheet", "document", "text", "audio", "unknown"
    extension: str          # e.g. ".png"
    strategy: str           # e.g. "image_multimodal", "image_text_fallback", "text_injected", ...
    text_content: str | None = None     # text injected into prompt (for non-image files)
    image_data_uri: str | None = None   # base64 data URI (for images in multimodal mode)
    success: bool = True
    error: str | None = None

    def to_dict(self) -> dict[str, Any]:
        return {
            "file_name": self.file_name,
            "file_type": self.file_type,
            "extension": self.extension,
            "strategy": self.strategy,
            "has_text_content": self.text_content is not None,
            "text_content_length": len(self.text_content) if self.text_content else 0,
            "has_image_data": self.image_data_uri is not None,
            "success": self.success,
            "error": self.error,
        }


@dataclass
class MultimodalContext:
    """Aggregated multimodal context for a single problem."""
    text_parts: list[str] = field(default_factory=list)
    image_data_uris: list[str] = field(default_factory=list)
    attachments: list[AttachmentResult] = field(default_factory=list)

    @property
    def has_images(self) -> bool:
        return len(self.image_data_uris) > 0

    @property
    def n_images(self) -> int:
        return len(self.image_data_uris)

    @property
    def extra_text(self) -> str:
        """Concatenated text from all attachments (for prompt injection)."""
        return "\n\n".join(self.text_parts)

    def summary_label(self) -> str:
        """Short human-readable label for logging."""
        parts = []
        for a in self.attachments:
            parts.append(f"{a.file_name}→{a.strategy}")
        return ", ".join(parts) if parts else "no_attachments"

    def to_metadata(self) -> list[dict[str, Any]]:
        return [a.to_dict() for a in self.attachments]


# ------------------------------------------------------------------
# File readers
# ------------------------------------------------------------------

def _encode_image_base64(path: str) -> str:
    """Read an image file and return a ``data:`` URI with base64 content."""
    mime = mimetypes.guess_type(path)[0] or "image/png"
    with open(path, "rb") as f:
        data = base64.b64encode(f.read()).decode("utf-8")
    return f"data:{mime};base64,{data}"


def _read_text_file(path: str, max_chars: int = MAX_TEXT_CHARS) -> str | None:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read(max_chars)
    except (IOError, OSError):
        return None


def _read_spreadsheet(path: str) -> str | None:
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".csv":
            return _read_text_file(path)
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        lines: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"=== Sheet: {sheet_name} ===")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                lines.append("\t".join(str(c) if c is not None else "" for c in row))
                row_count += 1
                if row_count > 2000:
                    lines.append("... (truncated)")
                    break
        wb.close()
        return "\n".join(lines)
    except ImportError:
        logger.warning("openpyxl not installed — cannot read .xlsx files; pip install openpyxl")
        return None
    except Exception as exc:
        logger.debug("Failed to read spreadsheet %s: %s", path, exc)
        return None


def _read_document(path: str) -> str | None:
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            return _read_pdf(path)
        if ext == ".docx":
            return _read_docx(path)
        if ext == ".pptx":
            return _read_pptx(path)
    except Exception as exc:
        logger.debug("Failed to read document %s: %s", path, exc)
    return None


def _read_pdf(path: str) -> str | None:
    for reader_fn in [_read_pdf_pdfplumber, _read_pdf_pypdf]:
        result = reader_fn(path)
        if result:
            return result
    logger.warning("No PDF reader available (install pdfplumber or PyPDF2)")
    return None


def _read_pdf_pdfplumber(path: str) -> str | None:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            pages = [p.extract_text() or "" for p in pdf.pages]
        text = "\n\n".join(pages)
        return text[:MAX_TEXT_CHARS] if text.strip() else None
    except ImportError:
        return None


def _read_pdf_pypdf(path: str) -> str | None:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(path)
        pages = [p.extract_text() or "" for p in reader.pages]
        text = "\n\n".join(pages)
        return text[:MAX_TEXT_CHARS] if text.strip() else None
    except ImportError:
        return None


def _read_docx(path: str) -> str | None:
    try:
        import docx
        doc = docx.Document(path)
        text = "\n".join(p.text for p in doc.paragraphs)
        return text[:MAX_TEXT_CHARS] if text.strip() else None
    except ImportError:
        logger.warning("python-docx not installed — cannot read .docx files")
        return None


def _read_pptx(path: str) -> str | None:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)
        text = "\n".join(texts)
        return text[:MAX_TEXT_CHARS] if text.strip() else None
    except ImportError:
        logger.warning("python-pptx not installed — cannot read .pptx files")
        return None


# ------------------------------------------------------------------
# Core: process a single file attachment
# ------------------------------------------------------------------

def process_attachment(
    file_path: str,
    file_name: str | None = None,
    multimodal: bool = True,
) -> AttachmentResult:
    """Process one file attachment and return an ``AttachmentResult``.

    Args:
        file_path: Absolute path to the file on disk.
        file_name: Human-readable file name (defaults to basename).
        multimodal: If True and file is an image, encode as base64 data URI.
                    If False, images produce a text placeholder instead.
    """
    fname = file_name or os.path.basename(file_path)
    ext = os.path.splitext(file_path)[1].lower()

    if not os.path.exists(file_path):
        return AttachmentResult(
            file_path=file_path, file_name=fname, file_type="missing",
            extension=ext, strategy="file_not_found", success=False,
            error=f"File does not exist: {file_path}",
        )

    # --- Images ---
    if ext in IMAGE_EXTENSIONS:
        if multimodal:
            try:
                data_uri = _encode_image_base64(file_path)
                return AttachmentResult(
                    file_path=file_path, file_name=fname, file_type="image",
                    extension=ext, strategy="image_multimodal",
                    image_data_uri=data_uri,
                )
            except Exception as exc:
                return AttachmentResult(
                    file_path=file_path, file_name=fname, file_type="image",
                    extension=ext, strategy="image_encode_failed",
                    text_content=f"[Image file: {fname} — failed to encode: {exc}]",
                    success=False, error=str(exc),
                )
        else:
            return AttachmentResult(
                file_path=file_path, file_name=fname, file_type="image",
                extension=ext, strategy="image_text_fallback",
                text_content=f"[Image file: {fname} — image content not available in text-only mode]",
            )

    # --- Text files ---
    if ext in TEXT_EXTENSIONS:
        content = _read_text_file(file_path)
        if content:
            return AttachmentResult(
                file_path=file_path, file_name=fname, file_type="text",
                extension=ext, strategy="text_injected",
                text_content=f"[Content of {fname}]:\n{content}",
            )
        return AttachmentResult(
            file_path=file_path, file_name=fname, file_type="text",
            extension=ext, strategy="text_read_failed",
            text_content=f"[Could not read file: {fname}]",
            success=False, error="Read failed",
        )

    # --- Spreadsheets ---
    if ext in SPREADSHEET_EXTENSIONS or ext == ".csv":
        content = _read_spreadsheet(file_path)
        if content:
            return AttachmentResult(
                file_path=file_path, file_name=fname, file_type="spreadsheet",
                extension=ext, strategy="spreadsheet_injected",
                text_content=f"[Spreadsheet content of {fname}]:\n{content}",
            )
        return AttachmentResult(
            file_path=file_path, file_name=fname, file_type="spreadsheet",
            extension=ext, strategy="spreadsheet_read_failed",
            text_content=f"[Could not read spreadsheet: {fname}]",
            success=False, error="Spreadsheet read failed",
        )

    # --- Documents (PDF, DOCX, PPTX) ---
    if ext in DOCUMENT_EXTENSIONS:
        content = _read_document(file_path)
        if content:
            return AttachmentResult(
                file_path=file_path, file_name=fname, file_type="document",
                extension=ext, strategy="document_injected",
                text_content=f"[Document content of {fname}]:\n{content}",
            )
        return AttachmentResult(
            file_path=file_path, file_name=fname, file_type="document",
            extension=ext, strategy="document_read_failed",
            text_content=f"[Could not extract text from: {fname}]",
            success=False, error="Document extraction failed",
        )

    # --- Unsupported / unknown ---
    return AttachmentResult(
        file_path=file_path, file_name=fname, file_type="unsupported",
        extension=ext, strategy="unsupported_type",
        text_content=f"[Attached file: {fname} (type {ext} — content not extracted)]",
    )


# ------------------------------------------------------------------
# Core: build multimodal context from one or more attachments
# ------------------------------------------------------------------

def build_multimodal_context(
    file_paths: list[tuple[str, str | None]],
    multimodal: bool = True,
) -> MultimodalContext:
    """Process a list of ``(file_path, file_name)`` pairs and return a
    :class:`MultimodalContext` that aggregates text and image content.

    Args:
        file_paths: List of ``(absolute_path, display_name_or_None)`` tuples.
        multimodal: Whether images should be encoded as base64 data URIs.
    """
    ctx = MultimodalContext()
    for fpath, fname in file_paths:
        result = process_attachment(fpath, fname, multimodal=multimodal)
        ctx.attachments.append(result)

        if result.text_content:
            ctx.text_parts.append(result.text_content)
        if result.image_data_uri:
            ctx.image_data_uris.append(result.image_data_uri)

        logger.info(
            "  Attachment: %s → %s%s",
            result.file_name,
            result.strategy,
            f" (error: {result.error})" if result.error else "",
        )

    return ctx


# ------------------------------------------------------------------
# Core: build OpenAI-compatible messages
# ------------------------------------------------------------------

def build_messages(
    system_prompt: str,
    user_text: str,
    image_data_uris: list[str] | None = None,
) -> list[dict[str, Any]]:
    """Build an OpenAI-compatible message array.

    When *image_data_uris* is provided, the user message uses the
    ``content`` list format (text + image_url entries).  Otherwise a
    plain string ``content`` is used.
    """
    if image_data_uris:
        content: list[dict[str, Any]] = [{"type": "text", "text": user_text}]
        for uri in image_data_uris:
            content.append({"type": "image_url", "image_url": {"url": uri}})
        return [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": content},
        ]

    return [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_text},
    ]


# ------------------------------------------------------------------
# Convenience: log a summary of multimodal processing for a problem
# ------------------------------------------------------------------

def log_multimodal_summary(
    problem_label: str,
    ctx: MultimodalContext,
    logger_instance: logging.Logger | None = None,
) -> None:
    """Log a concise summary of what multimodal content was used."""
    log = logger_instance or logger
    if not ctx.attachments:
        return
    for a in ctx.attachments:
        tag = "MULTIMODAL" if a.image_data_uri else "TEXT"
        status = "OK" if a.success else "FAIL"
        log.info(
            "  [%s] %s | %s | strategy=%s | %s",
            problem_label, tag, a.file_name, a.strategy, status,
        )
